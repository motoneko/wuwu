# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BLUE = RGBColor(0x42, 0x85, 0xF4)
RED = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xFB, 0xBC, 0x05)
GREEN = RGBColor(0x34, 0xA8, 0x53)
DARK = RGBColor(0x20, 0x21, 0x24)
GRAY = RGBColor(0x5F, 0x63, 0x68)
LIGHT = RGBColor(0xF1, 0xF3, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF8, 0xF9, 0xFA)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW = 13.333


def set_font(run, size=14, bold=False, color=DARK, font=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        set_font(run, size=size, bold=bold, color=color)
    return tb


def header(slide, num, title, subtitle=None, accent=BLUE):
    add_rect(slide, 0, 0, SW, 0.12, accent)
    add_rect(slide, 0.55, 0.42, 0.10, 0.62, accent)
    add_text(slide, 0.85, 0.38, 1.5, 0.6, num, size=26, bold=True, color=accent)
    add_text(slide, 2.05, 0.38, 10.6, 0.6, title, size=23, bold=True, color=DARK)
    if subtitle:
        add_text(slide, 2.05, 1.02, 10.6, 0.4, subtitle, size=12, color=GRAY)


def add_table(slide, x, y, w, rows_data, col_widths=None, header_color=BLUE,
              font_size=10.5, header_font_size=11, row_h=0.34, header_h=0.38,
              bold_first_col=False):
    nrows = len(rows_data); ncols = len(rows_data[0])
    total_h = header_h + row_h * (nrows - 1)
    gfx = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(total_h))
    table = gfx.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Emu(int(Inches(w) * cw / total))
    table.rows[0].height = Inches(header_h)
    for ri in range(1, nrows):
        table.rows[ri].height = Inches(row_h)
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = ROW_ALT
            else:
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            for li, ln in enumerate(str(val).split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = ln
                if ri == 0:
                    set_font(run, size=header_font_size, bold=True, color=WHITE)
                else:
                    set_font(run, size=font_size, bold=(bold_first_col and ci == 0), color=DARK)
    return table


def callout(slide, x, y, w, h, text, accent=YELLOW, size=12, bold=False):
    add_rect(slide, x, y, 0.08, h, accent)
    box = add_rect(slide, x + 0.08, y, w - 0.08, h, LIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.12)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        set_font(run, size=size, bold=bold, color=DARK)
    return box


colors = [BLUE, RED, YELLOW, GREEN]

# ============ Slide 1: Cover ============
s = add_slide()
add_rect(s, 0, 0, SW, 7.5, DARK)
for i, c in enumerate(colors):
    add_rect(s, 0.55 + i * 0.55, 1.95, 0.42, 0.12, c)
add_text(s, 0.55, 2.35, 12.2, 1.0, "Pixel 影像进化方向与迭代策略", size=40, bold=True, color=WHITE)
add_text(s, 0.55, 3.42, 12.2, 0.5, "交叉验证修订版｜整机产品经理视角｜Pixel 1–10 全系硬件策略与「硬件 for 软件」论证",
         size=17, color=RGBColor(0xBD, 0xC1, 0xC6))
add_text(s, 0.55, 4.35, 12.2, 0.5, "本版变化：经多来源交叉验证，修正 3 条原结论、新增 3 条结论（▲修正 ●新增）",
         size=14, bold=True, color=YELLOW)
add_text(s, 0.55, 5.85, 12.2, 1.1,
         "资料来源：Google 官方博客 / Google Research 技术博客 / Google I/O / Made by Google 发布会（2023–2025）\n核心人物访谈（Marc Levoy、Isaac Reynolds）/ DXOMARK / DPReview / Washington Post / PhoneArena\nAndroid Authority / TechRadar / Tom's Guide / Counterpoint 市场数据",
         size=11, color=GRAY, line_spacing=1.3)

# ============ Slide 2: TOC ============
s = add_slide()
header(s, "目录", "Pixel 影像进化策略报告（修订版）", accent=BLUE)
toc = [
    ("01", "Pixel 1–10 影像硬件总览矩阵", BLUE),
    ("02", "Pixel 1–7 阶段策略复盘", BLUE),
    ("03", "Pixel 8 系列影像（补齐）", BLUE),
    ("04", "Pixel 9 发布会硬件宣传方式（补齐）", BLUE),
    ("05", "Pixel 10 发布会硬件宣传方式（补齐）", BLUE),
    ("06", "Pixel 9 vs 10 硬件参数对比", BLUE),
    ("07", "四阶段进化模型（含组织证据）", GREEN),
    ("08", "核心决策与权衡判断", GREEN),
    ("09", "核心结论（交叉验证修订版 8 条）▲●", RED),
    ("10", "影像调教哲学 ●新增", RED),
    ("11", "硬件 for 软件｜三层拆解 ▲修正", RED),
    ("12", "硬件 ↔ AI 功能支撑映射", YELLOW),
    ("13", "Tensor ↔ 影像 AI 演进对照", YELLOW),
    ("14", "第三方交叉验证：DXOMARK ●新增", RED),
    ("15", "未来路径预测 ＆ 一句话总结", GREEN),
]
col_x = [0.85, 7.0]
for i, (num, t, c) in enumerate(toc):
    cx = col_x[i // 8]
    cy = 1.5 + (i % 8) * 0.7
    add_rect(s, cx, cy + 0.05, 0.09, 0.4, c)
    add_text(s, cx + 0.25, cy, 0.7, 0.5, num, size=15, bold=True, color=c)
    add_text(s, cx + 0.95, cy + 0.02, 5.0, 0.5, t, size=13.5, color=DARK)

# ============ Slide 3: 01 标准版矩阵 ============
s = add_slide()
header(s, "01", "Pixel 1–10 影像硬件总览矩阵 ─ 标准版", "十年硬件演进一图速览")
rows = [
    ["代次", "主摄", "超广角", "长焦"],
    ["Pixel 1–3", "12.2MP 单摄", "—", "—"],
    ["Pixel 4 / 4 XL", "12.2MP", "—", "16MP, 1.0µm, f/2.4, ~48mm（≈1.7x）"],
    ["Pixel 5", "12.2MP", "16MP, 1.0µm, f/2.2, 107°", "—"],
    ["Pixel 6", "50MP, 1/1.31\"", "12MP, 1.25µm, f/2.2, 114°", "—"],
    ["Pixel 7–8", "50MP, 1/1.31\"", "12MP, 1/2.9\", 1.25µm", "—"],
    ["Pixel 9", "50MP, 1/1.31\"", "48MP, f/1.7, 123°, 1/2.55\", PDAF", "—"],
    ["Pixel 10", "48MP, 1/2\"  ▼降档", "13MP, f/2.2, 120°, 1/3.1\"  ▼降档", "10.8MP, 1/3.2\", f/3.1, 5x 光学  ★新增"],
]
add_table(s, 0.55, 1.6, 12.2, rows, col_widths=[2.0, 3.2, 3.5, 3.5], row_h=0.5, font_size=11.5, header_font_size=12.5, bold_first_col=True)
callout(s, 0.55, 6.0, 12.2, 0.85,
        "★ Pixel 10 标准版：用 5x 长焦下放换主摄 / 超广降档 →「焦段完整性优先于单颗规格」",
        accent=RED, size=13, bold=True)

# ============ Slide 4: 01 Pro 矩阵 ============
s = add_slide()
header(s, "01", "Pixel 1–10 影像硬件总览矩阵 ─ Pro 版", "Pixel 6 起 Pro 三摄结构 4 代未变")
rows = [
    ["代次", "主摄", "超广角", "长焦"],
    ["Pixel 6 Pro", "50MP, 1/1.31\"", "12MP, f/2.2, 114°, 1.25µm", "48MP, 1/2\", f/3.5, 4x 光学"],
    ["Pixel 7 Pro", "50MP, 1/1.31\"", "12MP, f/2.2, 125.8°, 1/2.9\"", "48MP, 1/2.55\", OIS, 5x 光学"],
    ["Pixel 8 Pro", "50MP, 1/1.31\"（GN2）", "48MP, f/1.95, 1/2.0\", 125.5°", "48MP, f/2.8, 1/2.55\", PDAF, OIS, 5x"],
    ["Pixel 9 Pro / XL", "50MP, 1/1.31\"", "48MP, f/1.7, 1/2.55\", 123°", "48MP, 1/2.55\", OIS, 5x 光学"],
    ["Pixel 10 Pro / XL", "50MP, 1/1.3\"", "48MP, f/1.7, 1/2.55\", 123°", "48MP, f/2.8, 113mm, 5x, OIS"],
]
add_table(s, 0.55, 1.6, 12.2, rows, col_widths=[2.2, 3.2, 3.3, 3.5], row_h=0.55, font_size=11.5, header_font_size=12.5, bold_first_col=True)
callout(s, 0.55, 5.6, 12.2, 1.1,
        "★ Pixel 8 Pro = 硬件最后一次大改（三摄全 48MP+ 级、光量大升）\n★ Pixel 9 Pro → 10 Pro 硬件几乎不变，变焦却从 30x → 100x ─ 差异化已交给 AI",
        accent=BLUE, size=13, bold=True)

# ============ Slide 5: 02 复盘 ============
s = add_slide()
header(s, "02", "Pixel 1–7 阶段策略复盘", "动作 → 逻辑 → 结果（含商业与人事证据）")
rows = [
    ["代次", "硬件动作", "背后逻辑", "结果"],
    ["Pixel 1–3", "12MP 单摄，无超广无长焦", "「计算摄影」：单摄+AI 证明软件可重构手机摄影", "DXO 击败传统双/三/四摄旗舰"],
    ["Pixel 4 / 4 XL", "首次单→双：加 16MP 1.7x 长焦", "强化「Pixel 化」场景：变焦、人像、远景裁切", "焦段落后被吐槽；商业失败：两季度仅约 200 万台"],
    ["Pixel 5", "取消长焦、加 16MP 107° 超广", "2020 年超广已成「主流必需」", "路线摇摆；同年 3 月计算摄影负责人 Marc Levoy 离职"],
    ["Pixel 6 ★拐点", "自研 Tensor G1；50MP 平台重启；Pro 三摄", "① 12.2MP 算法到顶 ② Tensor 需要影像平台 ③ 对商业失败的纠错", "硬件大升级起点，奠定 50MP 黄金平台"],
    ["Pixel 7", "Pro 长焦 4x→5x，1/2.55\"+OIS", "服务远摄与 10x 裁切", "进入「硬件稳定、算法迭代」节奏"],
    ["Pixel 8", "主摄换 GN2（+35% 光）；Pro 全模组升级；标准版超广加 AF", "给生成式 AI 编辑功能更好的「原材料」", "Best Take / Magic Editor / Video Boost 集中爆发"],
]
add_table(s, 0.45, 1.55, 12.45, rows, col_widths=[1.7, 3.4, 3.8, 3.4], row_h=0.78, font_size=10.5, header_font_size=12, bold_first_col=True)

# ============ Slide 6: 03 P8 硬件 ============
s = add_slide()
header(s, "03", "Pixel 8 系列影像 ─ 核心硬件（补齐）", "Made by Google 2023.10.04 发布", accent=BLUE)
add_text(s, 0.55, 1.45, 6.0, 0.4, "Pixel 8（标准版）", size=15, bold=True, color=BLUE)
rows8 = [
    ["模组", "配置"],
    ["主摄", "50MP OctaPD, f/1.68, 1/1.31\"（GN2，光量+21%）"],
    ["超广角", "12MP, f/2.2, 125.8°, 1/2.9\"\n★首次加 AF → Macro Focus 下放"],
    ["前摄", "10.5MP DualPD, AF, 95°"],
    ["视频", "4K 24/30/60"],
]
add_table(s, 0.55, 1.9, 6.0, rows8, col_widths=[1.4, 4.6], row_h=0.62, font_size=10.5, bold_first_col=True)
add_text(s, 6.95, 1.45, 6.0, 0.4, "Pixel 8 Pro", size=15, bold=True, color=RED)
rows8p = [
    ["模组", "配置｜同代变化"],
    ["主摄", "50MP, f/1.68, 1/1.31\"｜换 Samsung GN2，光量+35%"],
    ["超广角", "48MP, f/1.95, 1/2.0\"｜光量+105%，对焦-30%"],
    ["长焦", "48MP, f/2.8, 5x, OIS, PDAF｜12→48MP，光量+56%"],
    ["前摄", "10.5MP DualPD｜首次加 AF"],
]
add_table(s, 6.95, 1.9, 5.85, rows8p, col_widths=[1.3, 4.55], row_h=0.62, font_size=10.5, header_color=RED, bold_first_col=True)
callout(s, 0.55, 5.7, 12.25, 1.2,
        "系列定位：「硬件给生成式 AI 喂素材」的集中爆发期\n所有硬件升级用「光量提升 %」表达 ─ 给 Best Take / Magic Editor / Video Boost 更好的原材料",
        accent=RED, size=12.5, bold=True)

# ============ Slide 7: 03 P8 宣传 ============
s = add_slide()
header(s, "03", "Pixel 8 系列 ─ 官网与发布会宣传方式（补齐）", "store.google.com + Made by Google 2023", accent=BLUE)
rows = [
    ["维度", "Pixel 8", "Pixel 8 Pro"],
    ["官网主标题", "\"The Google AI phone with the best\nPixel Camera yet\"", "\"The all-pro Google phone.\"\n\"Award-winning camera. Pro-level photos.\""],
    ["副 tagline", "\"Super helpful\" AI 助手叙事", "Pro Controls（手动 ISO/快门/白平衡 + 50MP RAW）"],
    ["影像卖点", "Best Take / Magic Editor / Macro Focus /\nAudio Magic Eraser", "＋Video Boost / Night Sight Video /\nPro Controls / Zoom Enhance（Pro 独享）"],
]
add_table(s, 0.55, 1.55, 12.2, rows, col_widths=[1.8, 5.2, 5.2], row_h=0.85, font_size=10.5, header_color=BLUE, bold_first_col=True)
rows2 = [
    ["发布会维度", "宣传方式"],
    ["整体叙事", "「Built Different」+ NBA 合作（Jimmy Butler、Giannis）；AI 是核心卖点"],
    ["硬件话术", "不堆参数，用「光量提升 %」表达硬件意义：主摄 +35%、超广 +105%、长焦 +56%"],
    ["软硬耦合", "Tensor G3 = \"AI-first processor\"，\"paving the way for on-device generative AI\""],
    ["差异化硬件", "温度传感器（Pro 独享）首次出现在智能手机，强化 Pro 身份"],
]
add_table(s, 0.55, 4.45, 12.2, rows2, col_widths=[1.8, 10.4], row_h=0.6, font_size=10.5, header_color=GRAY, bold_first_col=True)

# ============ Slide 8: 04 P9 发布会 ============
s = add_slide()
header(s, "04", "Pixel 9 系列 ─ 发布会硬件宣传方式（补齐）", "Made by Google 2024.08.13", accent=BLUE)
rows = [
    ["维度", "宣传方式"],
    ["核心叙事", "「Pixel 是 Google AI 的硬件入口」：Tensor G4 / Gemini Live / Pixel Studio / Pixel Screenshots"],
    ["硬件话术", "硬件参数让位 AI 叙事：① 相机条容纳更大传感器 ② Pro 三摄完整焦段 ③ 标准版超广 12→48MP"],
    ["软硬耦合", "Tensor G4 = \"first processor running Gemini Nano with multimodality\"（多模态是影像 AI 的前提）"],
    ["产品线动作", "Pro 拆分 6.3\"（Pro）/ 6.8\"（Pro XL），小尺寸首次获得完整 Pro 影像"],
    ["重点 AI 功能", "Add Me / Pixel Studio / Magic Editor-Reimagine / Zoom Enhance（Pro 独享 30x）"],
    ["视频亮点", "SuperRes Zoom Video 20x（Pro）/ Night Sight Video / 8K（Video Boost 云端）"],
]
add_table(s, 0.55, 1.55, 12.2, rows, col_widths=[2.0, 10.2], row_h=0.66, font_size=10.5, header_color=BLUE, bold_first_col=True)
callout(s, 0.55, 6.1, 12.2, 1.0,
        "系列总结：从「硬件叙事」转向「AI 叙事」的转折代 ─ 硬件提供稳定输入端，差异化由 AI / Gemini 拉开\n定位：用更强硬件 + AI 编辑解决「拍不好、拍不全、拍完想补救」",
        accent=BLUE, size=12, bold=True)

# ============ Slide 9: 05 P10 发布会 ============
s = add_slide()
header(s, "05", "Pixel 10 系列 ─ 发布会硬件宣传方式（补齐）", "Made by Google 2025.08.20, Brooklyn｜Jimmy Fallon 主持，Stephen Curry / Lando Norris 站台", accent=BLUE)
rows = [
    ["维度", "宣传方式"],
    ["核心叙事", "Tensor G5 + Gemini Nano = 拍摄全流程：拍前构图（Camera Coach）/ 拍中变焦（Pro Res Zoom）/ 拍后自然语言编辑（Ask Photos）"],
    ["芯片话术", "Tensor 五年最大升级：首次 TSMC 3nm（N3E）/ TPU +60% / CPU +34% / 全新 ISP（10-bit + motion deblur 默认开启）"],
    ["传感器话术", "标准版关键升级 = 5x 长焦下放；主摄 50→48MP、超广 48→13MP ─ 预算让位给长焦 + Tensor G5"],
    ["核心爆点", "100x Pro Res Zoom（Pro 独享）：G5 全新 ISP 跑 single-step diffusion，完全 on-device 秒级出图\n\"the largest model ever to run in Pixel Camera\""],
    ["重点 AI 功能", "Camera Coach（需联网）/ Auto Best Take（150 帧）/ Ask Photos / Add Me 升级"],
    ["真实性叙事", "C2PA 凭证全流程 on-device（G5 + Titan M2）；Pro Res Zoom 检测到人脸即禁用生成式 AI（guardrail）"],
]
add_table(s, 0.55, 1.65, 12.2, rows, col_widths=[2.0, 10.2], row_h=0.72, font_size=10.5, header_color=BLUE, bold_first_col=True)
callout(s, 0.55, 6.55, 12.2, 0.75,
        "系列总结：「拍后修」→「全流程介入」的范式转折，「硬件为 AI 输入端」最彻底的一代",
        accent=RED, size=12.5, bold=True)

# ============ Slide 10: 06 9v10 标准版 ============
s = add_slide()
header(s, "06", "Pixel 9 vs Pixel 10 ─ 标准版对比", "标准版最大变化：长焦补齐，主摄/超广让位")
rows = [
    ["模组", "Pixel 9", "Pixel 10", "变化判断"],
    ["后摄系统", "双摄（50 主 + 48 超广）", "三摄（48 主 + 13 超广 + 10.8 5x 长焦）", "★最大升级：新增 5x 长焦"],
    ["主摄", "50MP OctaPD, f/1.68, 1/1.31\"", "48MP QuadPD, f/1.70, 1/2\"", "▼传感器规格反而更小"],
    ["超广", "48MP QuadPD, f/1.7, 1/2.55\", AF", "13MP, f/2.2, 120°, 1/3.1\"", "▼明显降档"],
    ["长焦", "无；SuperRes Zoom 8x", "10.8MP 5x, OIS, SuperRes Zoom 20x", "★焦段补齐"],
    ["前摄", "10.5MP DualPD, AF, 95°", "同左", "延续"],
    ["视频变焦", "数字 7x", "数字 20x", "远摄视频增强"],
]
add_table(s, 0.55, 1.6, 12.2, rows, col_widths=[1.6, 3.6, 3.9, 3.1], row_h=0.6, font_size=11, bold_first_col=True)
callout(s, 0.55, 6.1, 12.2, 0.8,
        "判断：「降两颗、加一颗」是主动 BOM 重分配 ─ 焦段完整性的用户感知 > 主摄底大小的参数感知",
        accent=BLUE, size=12.5, bold=True)

# ============ Slide 11: 06 9v10 Pro ============
s = add_slide()
header(s, "06", "Pixel 9 vs Pixel 10 ─ Pro / Pro XL 对比", "Pro 版硬件停滞，AI 把可宣传能力推到 100x")
rows = [
    ["模组", "Pixel 9 Pro / XL", "Pixel 10 Pro / XL", "变化判断"],
    ["后摄系统", "50 + 48 + 48（5x）", "50 + 48 + 48（5x）", "硬件架构延续"],
    ["主摄", "50MP, f/1.68, 1/1.31\"", "50MP, f/1.68, 1/1.3\"", "几乎不变"],
    ["超广", "48MP, f/1.7, 1/2.55\"", "同左", "不变"],
    ["长焦", "48MP 5x, f/2.8, 1/2.55\"", "同左 + OIS 增强", "重心转 ProZoom"],
    ["前摄", "42MP DualPD, AF, 103°", "同左", "不变"],
    ["照片变焦", "SuperRes Zoom 30x", "ProZoom 100x（diffusion + Tensor G5）", "★Pro 最大升级：AI 长焦链路"],
    ["视频", "8K VideoBoost、20x SuperRes Video", "同左 + 更强稳定性", "防抖宣传增强"],
]
add_table(s, 0.55, 1.6, 12.2, rows, col_widths=[1.6, 3.7, 3.8, 3.1], row_h=0.54, font_size=11, bold_first_col=True)
callout(s, 0.55, 6.3, 12.2, 0.75,
        "★ 核心范式：「硬件停滞」+「AI 把同一套硬件可宣传能力推高 3 倍（30x→100x）」",
        accent=GREEN, size=13, bold=True)

# ============ Slide 12: 07 四阶段模型 ============
s = add_slide()
header(s, "07", "PM 视角｜四阶段进化模型（含组织证据）", "阶段划分 + 商业与人事佐证", accent=GREEN)
stages = [
    ("阶段一", "Pixel 1–3\n2016–2018", "单摄证明\n「软件重构摄影」", BLUE),
    ("阶段二", "Pixel 4–5\n2019–2020", "硬件补焦段\n「方向摇摆」", YELLOW),
    ("阶段三", "Pixel 6–8\n2021–2023", "自研芯片重启\n「50MP 黄金平台」", RED),
    ("阶段四", "Pixel 9–10\n2024–2025", "AI 输入端\n「全流程介入」", GREEN),
]
bw = 2.75
for i, (t1, t2, t3, c) in enumerate(stages):
    x = 0.55 + i * (bw + 0.35)
    add_rect(s, x, 1.5, bw, 0.45, c)
    add_text(s, x, 1.56, bw, 0.4, t1, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, 1.95, bw, 1.25, LIGHT)
    add_text(s, x + 0.1, 2.05, bw - 0.2, 0.6, t2, size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 2.6, bw - 0.2, 0.7, t3, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, x + bw, 2.1, 0.38, 0.5, "▶", size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
rows = [
    ["阶段", "核心假设", "关键动作", "阶段产物"],
    ["① 计算摄影证明", "软件可以重构手机摄影", "单摄 + HDR+ + 多帧合成", "DXO 击败双/三/四摄旗舰"],
    ["② 硬件补焦段", "算法不能完全替代焦段", "4 加长焦、5 换超广，路线摇摆", "商业失败（P4 两季度约 200 万台）"],
    ["③ 自研芯片驱动", "Tensor 需要影像平台展示 AI/ISP", "50MP 黄金平台多代复用", "算法长期调教 + 生成式 AI 爆发"],
    ["④ 全流程 AI 介入", "Pixel = Google AI 硬件入口", "硬件趋稳，AI 推高可宣传上限", "ProZoom 100x / Camera Coach"],
]
add_table(s, 0.55, 3.5, 12.2, rows, col_widths=[2.2, 3.2, 3.6, 3.2], row_h=0.52, header_h=0.36, font_size=10, header_color=GREEN, bold_first_col=True)
callout(s, 0.55, 6.25, 12.2, 1.0,
        "●组织铁证：HDR+/Night Sight 之父 Marc Levoy 2020 年 3 月离职（后任 Adobe 副总裁），时点正在 Pixel 4 商业失败之后、\nPixel 6 立项之时 ─「纯软件路线」的终结不只是技术到顶，更是商业失败 + 团队重组；Pixel 6 硬件重启本质是纠错",
        accent=RED, size=11.5, bold=True)

# ============ Slide 13: 08 决策权衡 ============
s = add_slide()
header(s, "08", "PM 视角｜核心决策与权衡判断", "8 个关键决策点的取舍逻辑", accent=GREEN)
rows = [
    ["决策点", "取舍", "PM 判断"],
    ["Pixel 4 选长焦不选超广", "长焦 =「Pixel 化」人像/远景", "押错方向：超广 2020 已成必需，焦段受损"],
    ["Pixel 5 砍长焦换超广", "跟随主流必需配置", "修正方向，代价是路线摇摆"],
    ["Pixel 6 启动黄金平台", "50MP 1/1.31\" + Tensor G1", "★历史性决策：对商业失败的纠错 + 4 代复用底座"],
    ["Pixel 6–9 主摄不换底", "稳定 4 代便于长期算法调教", "与国产「逐代换大底」反向：Quad Bayer 平衡 HDR/暗光"],
    ["Pixel 8 全面升光量", "GN2 + 超广/长焦像素翻倍", "为生成式 AI 喂素材（Best Take / Magic Editor）"],
    ["Pixel 9 Pro 拆双尺寸", "6.3\" / 6.8\" 同一套 Pro 相机", "拆需求层；模组空间约束超广无法上大底"],
    ["Pixel 10 标准版降配增焦", "主摄/超广降档 + 5x 长焦", "★焦段完整性 > 单颗规格；也是对行业长焦下放趋势的跟进"],
    ["Pixel 10 Pro 硬件不动", "100x 全靠 Tensor G5 + diffusion", "★硬件趋稳 → AI 推高上限，黄金平台策略奏效"],
]
add_table(s, 0.55, 1.55, 12.2, rows, col_widths=[3.0, 3.6, 5.6], row_h=0.6, font_size=10.5, header_color=GREEN, bold_first_col=True)

# ============ Slide 14: 09 核心结论 1/2 ============
s = add_slide()
header(s, "09", "核心结论（交叉验证修订版）─ 上", "▲= 经交叉验证修正的结论", accent=RED)
concls = [
    ("①", "硬件优先级：主摄 > 长焦 > 超广（Pro 例外）", BLUE,
     "主摄定下限，长焦定 Pro 溢价与宣传上限，超广定焦段完整性；Pro 版把超广作为 Macro 与三摄完整体验的补足"),
    ("②▲", "黄金平台策略奏效 ─ 但有可测量的天花板", RED,
     "DXOMARK：Pixel 10 Pro XL 163 分仅全球第 4，落后华为 Pura 80 Ultra / Oppo Find X8 Ultra / vivo X200 Ultra（1 英寸大底阵营）。调教保第一梯队但拿不了第一 ─ Google 主动放弃画质榜军备竞赛，改打「AI 功能可感知度」"),
    ("③▲", "P10 标准版降档 = 预算让位；但长焦下放是行业趋同，非 Google 引领", RED,
     "iPhone 15 Pro Max 5x tetraprism（2023）、三星双长焦、2025 潜望下探中端 ─「长焦是未来方向」应修正为「长焦已是高端标配，P10 是补行业课 + 给 AI Zoom 喂光学起点」"),
    ("④", "Pro 硬件 4 代不变、AI 独立承担差异化 ─ 有第三方实证", GREEN,
     "DXOMARK 确认 10 Pro XL 与 9 Pro XL 相机模组完全相同，但成绩显著提升（归因 G5 新 ISP）─「同硬件、软件提分」获测试机构背书"),
]
y = 1.5
for num, title, c, body in concls:
    add_rect(s, 0.55, y, 0.09, 1.25, c)
    add_text(s, 0.85, y, 0.85, 0.5, num, size=18, bold=True, color=c)
    add_text(s, 1.75, y, 11.0, 0.4, title, size=14.5, bold=True, color=DARK)
    add_text(s, 1.75, y + 0.44, 11.0, 0.8, body, size=10.5, color=GRAY, line_spacing=1.15)
    y += 1.42

# ============ Slide 15: 09 核心结论 2/2 ============
s = add_slide()
header(s, "09", "核心结论（交叉验证修订版）─ 下", "●= 交叉验证后新增的结论", accent=RED)
concls = [
    ("⑤▲", "「硬件 for 软件」需分三层：营销层 / 产品层 / 公司层", RED,
     "营销层：硬件为 AI 功能服务（成立）；产品层：硬件投入压到支撑 AI 叙事的最低水平（全球份额 ~1%，BOM 克制）；公司层：Pixel = Google AI 的 halo device（北美份额 4.8%→12.8%，高端增速 105% YoY 全球第一）"),
    ("⑥●", "视频是「软件补硬件」的失效反例", YELLOW,
     "低光视频/动态范围多年落后 iPhone；Video Boost 云端处理最长近 20h/条，三代后仍被评「beta」─ 算法无法端侧实时完成时此路不通；G5 全新 ISP（10-bit 默认）是用真硬件还视频欠账"),
    ("⑦●", "调教北极星：\"They're memories, not photos\"", GREEN,
     "Pixel Camera 负责人 Isaac Reynolds：从「图像处理」转向「记忆重建」；The Pixel Look（暗部对比强、高光严格保护）是故意的品牌资产 ─ 所有软硬件投入最终服务「更自信的拍摄者 + 更美的记忆」"),
    ("⑧●", "真实性悖论是战略级风险，Google 预先把风险变卖点", BLUE,
     "Best Take 被批「保存从未发生的时刻」（Washington Post）；应对组合拳：人脸生成 guardrail + C2PA 全流程 on-device + \"true to your memory\" 话术 ─ 把行业争议预先转化为信任优势"),
]
y = 1.5
for num, title, c, body in concls:
    add_rect(s, 0.55, y, 0.09, 1.25, c)
    add_text(s, 0.85, y, 0.85, 0.5, num, size=18, bold=True, color=c)
    add_text(s, 1.75, y, 11.0, 0.4, title, size=14.5, bold=True, color=DARK)
    add_text(s, 1.75, y + 0.44, 11.0, 0.8, body, size=10.5, color=GRAY, line_spacing=1.15)
    y += 1.42

# ============ Slide 16: 10 调教哲学 ============
s = add_slide()
header(s, "10", "影像调教哲学（●新增）", "来自 Google Research 技术博客与核心人物访谈", accent=RED)
rows = [
    ["调教策略", "技术细节", "来源"],
    ["刻意欠曝", "HDR+ 按快门捕捉 3–15 帧全部欠曝，保高光不过曝", "Google Research（HDR+ with Bracketing）"],
    ["局部色调映射", "按区域亮度/纹理/噪声施加不同曲线，\"detail everywhere\"", "Google Research"],
    ["风格故意「浓郁」", "暗部深、对比强的 dramatic look 是设计选择：\"distinct and cool and fresh\"", "Isaac Reynolds 访谈"],
    ["包围曝光补暗部", "Pixel 5 起 HDR+ with Bracketing 解决欠曝带来的暗部噪声", "Google Research"],
    ["虚化超越物理镜头", "人像虚化「不受真实镜头限制」，比真实光学更可控", "Isaac Reynolds"],
    ["Real Tone 进 ISP", "与有色人种影像专家合作扩充训练集，调教面部检测/白平衡/曝光；G5 ISP 继续强化", "Google I/O 2021 / Image Equity"],
    ["人脸生成 guardrail", "Pro Res Zoom 30x 以上才启用生成式 AI，检测到人脸即退回传统超分", "DPReview 实测"],
]
add_table(s, 0.55, 1.55, 12.2, rows, col_widths=[2.4, 6.2, 3.6], row_h=0.56, font_size=10.5, header_color=RED, bold_first_col=True)
callout(s, 0.55, 6.15, 12.2, 1.0,
        "哲学三段论：拍到（capture, P1–5 Levoy 时代）→ 拍好（process, P6–9）→ 记忆（recreate, P10–）\nReynolds 内部使命：\"Create more confident photographers and more beautiful memories.\"",
        accent=GREEN, size=12, bold=True)

# ============ Slide 17: 11 三层拆解 ============
s = add_slide()
header(s, "11", "硬件 for 软件｜三层拆解（▲修正）", "不是单层结论：营销层 / 产品层 / 公司层 + 反例边界", accent=RED)
rows = [
    ["层级", "表述", "关键证据"],
    ["营销层", "硬件为 AI 功能服务（成立）", "每代硬件升级点对应 AI 功能落地点（见 12 节映射表）"],
    ["产品层", "硬件投入压到「支撑 AI 叙事的最低水平」", "全球份额 ~1% → BOM 克制、平台多代复用、预算转向自研芯片"],
    ["公司层", "Pixel = Google AI 的 halo device + 试验场", "北美份额 4.8%→12.8%（P9 后翻倍）；2025 高端增速 105% YoY 全球最快；技术反哺 Android 生态"],
]
add_table(s, 0.55, 1.55, 12.2, rows, col_widths=[1.6, 4.4, 6.2], row_h=0.72, font_size=11, header_color=RED, bold_first_col=True)
add_text(s, 0.55, 4.35, 12.2, 0.4, "反例与边界（●新增）", size=15, bold=True, color=DARK)
rows2 = [
    ["边界", "证据", "Google 的应对"],
    ["视频：「软件补硬件」失效", "Video Boost 云端最长 20h/条，三代仍像 beta；低光视频落后 iPhone", "G5 全新 ISP 用「真硬件」还欠账（10-bit 默认）"],
    ["画质榜：「调教补大底」失效", "DXOMARK 第 4，落后 3 家 1 英寸大底国产旗舰", "不堆硬件，改变战场 ─ 比「AI 功能可感知度」"],
]
add_table(s, 0.55, 4.8, 12.2, rows2, col_widths=[3.0, 5.2, 4.0], row_h=0.72, font_size=10.5, header_color=GRAY, bold_first_col=True)

# ============ Slide 18: 11b 三大叙事 ============
s = add_slide()
header(s, "11", "硬件 for 软件｜服务的三大软件叙事", "影像是 AI 叙事中最容易被用户感知的一环", accent=YELLOW)
cards = [
    ("叙事 A", "Google AI 个人助手", "Pixel = AI 入口", [
        "Gemini Nano on-device（多模态）", "Gemini Live 实时对话", "Pixel Studio 端侧生图",
        "Pixel Screenshots 截图理解", "Magic Cue / Voice Translate"], BLUE),
    ("叙事 B", "生成式 AI 影像", "Pixel = AI 影像创作平台", [
        "拍前：Camera Coach（Gemini 指导）", "拍中：Pro Res Zoom 100x（diffusion）",
        "拍中：Auto Best Take（150 帧分析）", "拍后：Magic Editor / Ask Photos",
        "拍后：Add Me / Zoom Enhance"], RED),
    ("叙事 C", "可信 / 包容性影像", "Pixel = 负责任的 AI", [
        "C2PA Content Credentials（on-device）", "Real Tone 多肤色准确还原",
        "人脸生成 guardrail", "Guided Frame 低视力辅助拍摄"], GREEN),
]
cw = 3.95
for i, (tag, title, sub, items, c) in enumerate(cards):
    x = 0.55 + i * (cw + 0.2)
    add_rect(s, x, 1.55, cw, 0.85, c)
    add_text(s, x + 0.15, 1.63, cw - 0.3, 0.35, tag + "｜" + title, size=14, bold=True, color=WHITE)
    add_text(s, x + 0.15, 2.02, cw - 0.3, 0.3, sub, size=11, color=WHITE)
    box = add_rect(s, x, 2.4, cw, 3.6, LIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.15)
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        run = p.add_run()
        run.text = "● " + it
        set_font(run, size=11, color=DARK)
callout(s, 0.55, 6.35, 12.2, 0.75,
        "硬件不只是为「手机上的软件」服务，更是为「证明 Google AI 领先」这个公司级叙事服务",
        accent=YELLOW, size=12.5, bold=True)

# ============ Slide 19: 12 硬件↔AI 映射 ============
s = add_slide()
header(s, "12", "硬件 for 软件｜核心支撑映射", "每一个硬件升级点 = 一个 AI 功能落地点（含官方原话）", accent=YELLOW)
rows = [
    ["硬件", "直接支撑的 AI / 软件功能", "论证（官方语言）"],
    ["Tensor G3 TPU（P8）", "Best Take / Magic Editor on-device 生成式 inpainting", "\"paving the way for on-device generative AI\""],
    ["Tensor G4 多模态（P9）", "Add Me / Pixel Studio / Zoom Enhance", "\"first processor running Gemini Nano with multimodality\""],
    ["Tensor G5 全新 ISP（P10）", "Pro Res Zoom 100x（diffusion on-device）/ 10-bit 视频", "\"the largest model ever to run in Pixel Camera\""],
    ["Tensor G5 + Titan M2", "C2PA Content Credentials 全流程 on-device", "\"full process takes place on-device\""],
    ["50MP Quad Bayer 主摄", "默认输出 12.5MP，四合一保留低光素材给 ML", "稳定平台上迭代计算摄影模型"],
    ["GN2 主摄 +35% 光（P8）", "Video Boost / Night Sight Video 的原材料", "\"pairs Tensor G3 with our powerful data centers\""],
    ["5x 长焦下放（P10 标准版）", "SuperRes Zoom 20x + AI Zoom 的光学起点", "生成式变焦需要光学起点，纯算法天花板低"],
    ["超广 AF（P8 起）", "Macro Focus 下放", "硬件 + 算法联合实现微距"],
    ["数据中心（云）", "Video Boost / Night Sight Video（云端补端侧）", "每分钟视频 = 1800 张照片量级的云端处理"],
]
add_table(s, 0.55, 1.55, 12.2, rows, col_widths=[2.9, 4.6, 4.7], row_h=0.56, font_size=10, header_color=YELLOW, header_font_size=11.5, bold_first_col=True)

# ============ Slide 20: 13 Tensor 对照 ============
s = add_slide()
header(s, "13", "Tensor ↔ 影像 AI 功能演进对照", "每一代芯片跃进 = 每一代 AI 功能跃迁", accent=YELLOW)
rows = [
    ["芯片", "机型 / 年份", "工艺与性能", "影像 AI 功能跃迁"],
    ["Tensor G1", "Pixel 6（2021）", "Samsung 5nm", "奠基：Real Tone / Magic Eraser / Motion Mode"],
    ["Tensor G2", "Pixel 7（2022）", "Samsung 5nm", "Photo Unblur / Cinematic Blur / Guided Frame"],
    ["Tensor G3", "Pixel 8（2023）", "Samsung 4nm｜\"AI-first processor\"", "★on-device 生成式 AI 起点：Best Take / Magic Editor / Video Boost（端云）"],
    ["Tensor G4", "Pixel 9（2024）", "Samsung 4nm｜Gemini Nano 多模态", "★Add Me / Pixel Studio / Reimagine / SuperRes Zoom 30x"],
    ["Tensor G5", "Pixel 10（2025）", "TSMC 3nm N3E｜TPU+60% CPU+34%｜全新 ISP", "★diffusion on-device：ProZoom 100x / Camera Coach / Auto Best Take / C2PA"],
]
add_table(s, 0.55, 1.55, 12.2, rows, col_widths=[1.6, 1.9, 3.6, 5.1], row_h=0.78, font_size=10.5, header_color=YELLOW, header_font_size=11.5, bold_first_col=True)
callout(s, 0.55, 6.1, 12.2, 1.0,
        "规律：芯片跃进 → AI 功能跃迁 → 传感器趋稳甚至降档（P10 主摄 1/2\"）\n硬件预算从「传感器」转移到「SoC / ISP / TPU / AI 模型部署」─「硬件为软件服务」最强力的实证",
        accent=RED, size=12.5, bold=True)

# ============ Slide 21: 14 DXOMARK 交叉验证 ============
s = add_slide()
header(s, "14", "第三方交叉验证：DXOMARK 数据（●新增）", "既是「硬件 for 软件」最硬的证据，也是它边界的最硬证据", accent=RED)
rows = [
    ["发现", "数据", "对结论的作用"],
    ["「同硬件、软件提分」实证", "10 Pro XL 与 9 Pro XL 相机模组完全相同，但总分/曝光/色彩/肤色显著提升（归因 G5 ISP 更细粒度对象分割）", "✔ 支撑「硬件 for 软件」与黄金平台策略"],
    ["画质榜天花板", "Pixel 10 Pro XL 163 分全球第 4，落后华为 Pura 80 Ultra / Oppo Find X8 Ultra / vivo X200 Ultra（1 英寸大底阵营）", "▲ 限定：调教补不齐大底差距，Google 主动不打这一仗"],
    ["9 代际提升点", "9 Pro XL vs 8 Pro：逆光保高光同时主体曝光准确，业界第一梯队", "✔ The Pixel Look 调教哲学的实测体现"],
    ["10 代际提升点", "10 Pro XL vs 9 Pro XL：肤色更自然（9 偏橙、10 中性）、白平衡更准", "✔ Real Tone 进 G5 ISP 流水线的实测体现"],
]
add_table(s, 0.55, 1.6, 12.2, rows, col_widths=[2.6, 5.8, 3.8], row_h=0.95, font_size=10.5, header_color=RED, bold_first_col=True)
callout(s, 0.55, 6.0, 12.2, 0.9,
        "PM 读法：支撑证据与天花板证据必须同时讲，结论才完整 ─\n「稳定硬件 + 调教」保第一梯队但拿不了第一；Google 的选择是改变战场，比「AI 功能可感知度」",
        accent=YELLOW, size=12, bold=True)

# ============ Slide 22: 15 未来路径 ============
s = add_slide()
header(s, "15", "未来路径预测（PM 视角）", "短期（Pixel 11/12）与中期（2 年外）", accent=GREEN)
add_text(s, 0.55, 1.5, 6.0, 0.4, "短期：Pixel 11 / 12", size=15, bold=True, color=BLUE)
rows1 = [
    ["维度", "预测"],
    ["传感器", "Pro 三件套（50+48+48 5x）几乎不变"],
    ["ISP", "G6 继续承担差异化；视频是最大欠账，on-device Video Boost 预计是下一个主打"],
    ["AI 功能", "「拍前」加重：Camera Coach 从建议走向实时取景介入"],
    ["标准版", "保留长焦；与 Pro 硬件差距收窄，靠 ProZoom / Pro Controls 拉差异"],
]
add_table(s, 0.55, 1.95, 6.0, rows1, col_widths=[1.3, 4.7], row_h=0.82, font_size=10.5, bold_first_col=True)
add_text(s, 6.95, 1.5, 6.0, 0.4, "中期：2 年外", size=15, bold=True, color=GREEN)
rows2 = [
    ["维度", "预测"],
    ["生成式成像", "100x 之后是「可控生成」：自然语言指导内容/风格/构图"],
    ["真实性博弈", "C2PA + 端侧凭证成行业标配，Pixel 提前布局获得话语权"],
    ["Pro 边界", "AI 弥补硬件差异后，「Pro」可能退回尺寸/续航/材质分层"],
]
add_table(s, 6.95, 1.95, 5.85, rows2, col_widths=[1.3, 4.55], row_h=1.0, font_size=10.5, header_color=GREEN, bold_first_col=True)

# ============ Slide 23: 金句 ============
s = add_slide()
add_rect(s, 0, 0, SW, 7.5, DARK)
for i, c in enumerate(colors):
    add_rect(s, 0.55 + i * 0.55, 1.5, 0.42, 0.12, c)
add_text(s, 0.55, 2.05, 12.2, 1.6, "「Pixel 不是一台拍照手机，是 Google AI 的硬件入口；\n  而影像的北极星不是画质榜，是『记忆』。」", size=28, bold=True, color=WHITE, line_spacing=1.35)
add_text(s, 0.55, 4.35, 12.2, 1.8,
         "主摄保下限，长焦撑溢价，超广补完整性；\n画质第一让给大底阵营，AI 可感知度自己拿走；\n硬件让位、AI 接管、真实性风险预先变卖点 ─ 这是 Pixel 10 完成的范式。",
         size=16, color=RGBColor(0xBD, 0xC1, 0xC6), line_spacing=1.5)

# ============ Slide 24: 附录 ============
s = add_slide()
header(s, "附录", "资料来源（多来源交叉论证）", "官方 / 技术博客 / 人物访谈 / 发布会 / 第三方测评与市场数据", accent=GRAY)
groups = [
    ("Google 官方 / 技术博客", [
        "blog.google – Tensor G5 / Tensor G3 / Pixel 8 Pro Camera",
        "Google Research – HDR+ with Bracketing / Live HDR+",
        "blog.google – Image Equity: Real Tone",
        "blog.google – Pro Res Zoom",
        "store.google.com – Pixel 8 Pro / Pixel 10 产品页",
    ], BLUE),
    ("人物访谈 / 发布会", [
        "DPReview – Marc Levoy 离职报道（2020.5）",
        "9to5Google – Isaac Reynolds「memories not photos」",
        "PetaPixel Podcast – Pixel 团队设计哲学（2025.8）",
        "TechCrunch / 9to5Google – MBG 2023",
        "Tom's Guide – MBG 2024；TechRadar – MBG 2025",
    ], RED),
    ("第三方测评 / 数据", [
        "DXOMARK – P8 Pro / P9 Pro XL / P10 Pro XL（163 分第 4）",
        "DPReview – Pro Res Zoom 实测（人脸 guardrail）",
        "Washington Post – Best Take 真实性争议",
        "PhoneArena / Android Authority – Video Boost 云端 20h",
        "Android Authority – Pixel 高端增速 105% YoY（H1 2025）",
        "PhoneArena – iPhone 15 Pro Max tetraprism（行业长焦）",
    ], GREEN),
]
cw = 3.95
for i, (title, items, c) in enumerate(groups):
    x = 0.55 + i * (cw + 0.2)
    add_rect(s, x, 1.6, cw, 0.5, c)
    add_text(s, x + 0.15, 1.69, cw - 0.3, 0.35, title, size=12.5, bold=True, color=WHITE)
    box = add_rect(s, x, 2.1, cw, 4.4, LIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.15)
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        run = p.add_run()
        run.text = "● " + it
        set_font(run, size=10, color=DARK)

out = "/mnt/c/Users/陈建壮/Desktop/code/wuwu/Pixel影像进化策略报告.pptx"
prs.save(out)
print(f"Saved: {out} | slides: {len(prs.slides._sldIdLst)}")
